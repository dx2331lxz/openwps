from __future__ import annotations

import base64
import hashlib
import io
import json
import logging
import re
import shutil
import subprocess
import uuid
from datetime import datetime, timezone
from pathlib import Path, PurePosixPath
from typing import Any

import httpx
from fastapi import HTTPException

from .config import BASE_DIR

logger = logging.getLogger("uvicorn.error")

WORKSPACES_DIR = BASE_DIR / "data" / "workspaces"
WORKSPACES_META_PATH = WORKSPACES_DIR / "meta.json"
DEFAULT_WORKSPACE_ID = "default"
FILES_DIRNAME = "files"
INTERNAL_DIRNAME = ".openwps"
REFERENCES_DIRNAME = "_references"
MEMORY_DIRNAME = "memory"
MEMORY_ENTRYPOINT_FILENAME = "MEMORY.md"
MEMORY_ROOT_RELATIVE = f"{INTERNAL_DIRNAME}/{MEMORY_DIRNAME}"
MEMORY_ENTRYPOINT_RELATIVE = f"{MEMORY_ROOT_RELATIVE}/{MEMORY_ENTRYPOINT_FILENAME}"
MEMORY_TYPES = {"user", "feedback", "project", "reference"}
MEMORY_ENTRYPOINT_MAX_LINES = 200
MEMORY_ENTRYPOINT_MAX_BYTES = 25_000
MEMORY_FILE_MAX_LINES = 200
MEMORY_FILE_MAX_BYTES = 4_096
MEMORY_SELECTED_LIMIT = 5
MEMORY_SCAN_LIMIT = 200
SYSTEM_ROOT_DIR_ORDER = {
    INTERNAL_DIRNAME: 0,
    REFERENCES_DIRNAME: 1,
}

ALLOWED_EXTENSIONS = {".docx", ".txt", ".md", ".pdf", ".ppt", ".pptx", ".markdown"}
EDITABLE_EXTENSIONS = {".docx", ".txt", ".md", ".markdown"}
READONLY_EXTENSIONS = ALLOWED_EXTENSIONS - EDITABLE_EXTENSIONS

MIME_MAP = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/vnd.ms-powerpoint": ".ppt",
    "application/pdf": ".pdf",
    "text/plain": ".txt",
    "text/markdown": ".md",
}

DEFAULT_PAGE_CONFIG: dict[str, int] = {
    "pageWidth": 794,
    "pageHeight": 1123,
    "marginTop": 96,
    "marginBottom": 96,
    "marginLeft": 113,
    "marginRight": 113,
}

DEFAULT_PARAGRAPH_ATTRS: dict[str, Any] = {
    "align": "left",
    "firstLineIndent": 0,
    "indent": 0,
    "rightIndent": 0,
    "headingLevel": None,
    "fontSizeHint": None,
    "fontFamilyHint": None,
    "lineHeight": 1.5,
    "spaceBefore": 0,
    "spaceAfter": 0,
    "listType": None,
    "listLevel": 0,
    "listChecked": False,
    "pageBreakBefore": False,
    "tabStops": [],
}

WORKSPACES_DIR.mkdir(parents=True, exist_ok=True)


def utc_now() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%S")


def _read_json(path: Path, fallback: Any) -> Any:
    if not path.exists():
        return fallback
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return fallback


def _write_json(path: Path, value: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2), encoding="utf-8")


def _slugify_workspace_id(value: str) -> str:
    text = re.sub(r"[^A-Za-z0-9_-]+", "-", value.strip()).strip("-_").lower()
    return text[:48] or uuid.uuid4().hex[:12]


def _workspace_root(workspace_id: str) -> Path:
    return WORKSPACES_DIR / workspace_id


def _workspace_files_root(workspace_id: str) -> Path:
    return _workspace_root(workspace_id) / FILES_DIRNAME


def _workspace_internal_root(workspace_id: str) -> Path:
    return _workspace_files_root(workspace_id) / INTERNAL_DIRNAME


def _workspace_memory_root(workspace_id: str) -> Path:
    return _workspace_internal_root(workspace_id) / MEMORY_DIRNAME


def _workspace_memory_entrypoint(workspace_id: str) -> Path:
    return _workspace_memory_root(workspace_id) / MEMORY_ENTRYPOINT_FILENAME


def _load_workspaces_meta() -> dict[str, Any]:
    raw = _read_json(WORKSPACES_META_PATH, {})
    if not isinstance(raw, dict):
        raw = {}
    items = raw.get("workspaces")
    if not isinstance(items, list):
        items = []
    return {
        "version": int(raw.get("version") or 1),
        "activeWorkspaceId": str(raw.get("activeWorkspaceId") or DEFAULT_WORKSPACE_ID),
        "workspaces": [item for item in items if isinstance(item, dict)],
    }


def _save_workspaces_meta(meta: dict[str, Any]) -> None:
    _write_json(WORKSPACES_META_PATH, meta)


def _workspace_entry(workspace_id: str, name: str | None = None) -> dict[str, Any]:
    now = utc_now()
    return {
        "id": workspace_id,
        "name": name or ("默认工作区" if workspace_id == DEFAULT_WORKSPACE_ID else workspace_id),
        "createdAt": now,
        "updatedAt": now,
    }


def _ensure_workspace_dirs(workspace_id: str) -> None:
    files_root = _workspace_files_root(workspace_id)
    (files_root / REFERENCES_DIRNAME).mkdir(parents=True, exist_ok=True)
    internal_root = files_root / INTERNAL_DIRNAME
    (internal_root / "index").mkdir(parents=True, exist_ok=True)
    (internal_root / "versions").mkdir(parents=True, exist_ok=True)
    (internal_root / "cache").mkdir(parents=True, exist_ok=True)
    memory_root = _workspace_memory_root(workspace_id)
    memory_root.mkdir(parents=True, exist_ok=True)
    entrypoint = _workspace_memory_entrypoint(workspace_id)
    if not entrypoint.exists():
        entrypoint.write_text("", encoding="utf-8")


def ensure_default_workspace() -> None:
    meta = _load_workspaces_meta()
    workspace_ids = {str(item.get("id") or "") for item in meta["workspaces"]}
    if DEFAULT_WORKSPACE_ID not in workspace_ids:
        meta["workspaces"].insert(0, _workspace_entry(DEFAULT_WORKSPACE_ID))
        meta["activeWorkspaceId"] = DEFAULT_WORKSPACE_ID
    _ensure_workspace_dirs(DEFAULT_WORKSPACE_ID)
    _save_workspaces_meta(meta)


def _require_workspace(workspace_id: str | None = None) -> str:
    ensure_default_workspace()
    target = str(workspace_id or get_active_workspace_id()).strip() or DEFAULT_WORKSPACE_ID
    meta = _load_workspaces_meta()
    if not any(str(item.get("id") or "") == target for item in meta["workspaces"]):
        raise HTTPException(status_code=404, detail="工作区不存在")
    _ensure_workspace_dirs(target)
    return target


def get_active_workspace_id() -> str:
    ensure_default_workspace()
    return _load_workspaces_meta().get("activeWorkspaceId") or DEFAULT_WORKSPACE_ID


def list_workspaces() -> dict[str, Any]:
    ensure_default_workspace()
    meta = _load_workspaces_meta()
    return {
        "activeWorkspaceId": meta.get("activeWorkspaceId") or DEFAULT_WORKSPACE_ID,
        "workspaces": meta.get("workspaces", []),
    }


def create_workspace(name: str | None = None, workspace_id: str | None = None) -> dict[str, Any]:
    ensure_default_workspace()
    meta = _load_workspaces_meta()
    base_id = _slugify_workspace_id(workspace_id or name or "workspace")
    existing = {str(item.get("id") or "") for item in meta["workspaces"]}
    candidate = base_id
    index = 2
    while candidate in existing:
        candidate = f"{base_id}-{index}"
        index += 1
    entry = _workspace_entry(candidate, name or candidate)
    meta["workspaces"].append(entry)
    meta["activeWorkspaceId"] = candidate
    _ensure_workspace_dirs(candidate)
    _save_workspaces_meta(meta)
    return entry


def set_active_workspace(workspace_id: str) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    meta = _load_workspaces_meta()
    meta["activeWorkspaceId"] = target
    _save_workspaces_meta(meta)
    return {"activeWorkspaceId": target}


def _sanitize_relative_path(raw_path: str | None, *, allow_empty: bool = False, allow_internal: bool = False) -> str:
    text = str(raw_path or "").replace("\\", "/").strip()
    if text in {"", "."}:
        if allow_empty:
            return ""
        raise HTTPException(status_code=400, detail="路径不能为空")
    pure = PurePosixPath(text)
    if pure.is_absolute():
        raise HTTPException(status_code=400, detail="路径必须是工作区相对路径")
    parts = [part for part in pure.parts if part not in {"", "."}]
    if any(part == ".." for part in parts):
        raise HTTPException(status_code=400, detail="路径不能包含 ..")
    if parts and parts[0] == INTERNAL_DIRNAME and not allow_internal:
        raise HTTPException(status_code=400, detail=".openwps 为系统目录，不能直接操作")
    if not parts:
        if allow_empty:
            return ""
        raise HTTPException(status_code=400, detail="路径不能为空")
    return "/".join(parts)


def _resolve_workspace_path(
    workspace_id: str,
    raw_path: str | None,
    *,
    allow_empty: bool = False,
    allow_internal: bool = False,
) -> tuple[str, Path]:
    relative = _sanitize_relative_path(raw_path, allow_empty=allow_empty, allow_internal=allow_internal)
    root = _workspace_files_root(workspace_id).resolve()
    target = (root / relative).resolve() if relative else root
    try:
        target.relative_to(root)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail="路径越界") from exc
    return relative, target


def _is_memory_workspace_path(raw_path: str | None) -> bool:
    text = str(raw_path or "").replace("\\", "/").strip().strip("/")
    return text == MEMORY_ROOT_RELATIVE or text.startswith(f"{MEMORY_ROOT_RELATIVE}/")


def _memory_workspace_path(memory_relative_path: str) -> str:
    clean = memory_relative_path.strip("/")
    return f"{MEMORY_ROOT_RELATIVE}/{clean}" if clean else MEMORY_ROOT_RELATIVE


def _sanitize_memory_relative_path(raw_path: str | None, *, allow_empty: bool = False) -> str:
    text = str(raw_path or "").replace("\\", "/").strip()
    if text.startswith(f"{MEMORY_ROOT_RELATIVE}/"):
        text = text[len(MEMORY_ROOT_RELATIVE) + 1:]
    elif text == MEMORY_ROOT_RELATIVE:
        text = ""
    if text in {"", "."}:
        if allow_empty:
            return ""
        raise HTTPException(status_code=400, detail="记忆文件路径不能为空")
    pure = PurePosixPath(text)
    if pure.is_absolute():
        raise HTTPException(status_code=400, detail="记忆文件路径必须是相对路径")
    parts = [part for part in pure.parts if part not in {"", "."}]
    if any(part == ".." for part in parts):
        raise HTTPException(status_code=400, detail="记忆文件路径不能包含 ..")
    if parts and parts[0] == INTERNAL_DIRNAME:
        raise HTTPException(status_code=400, detail="记忆文件路径必须位于 .openwps/memory 内")
    if not parts:
        if allow_empty:
            return ""
        raise HTTPException(status_code=400, detail="记忆文件路径不能为空")
    relative = "/".join(parts)
    suffix = Path(relative).suffix.lower()
    if suffix not in {".md", ".markdown"}:
        raise HTTPException(status_code=400, detail="记忆文件只支持 .md / .markdown")
    return relative


def _resolve_memory_path(workspace_id: str, raw_path: str | None) -> tuple[str, str, Path]:
    memory_relative = _sanitize_memory_relative_path(raw_path)
    root = _workspace_memory_root(workspace_id).resolve()
    target = (root / memory_relative).resolve()
    try:
        target.relative_to(root)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail="记忆文件路径越界") from exc
    return memory_relative, _memory_workspace_path(memory_relative), target


def _detect_extension(filename: str, content_type: str = "") -> str:
    ct = (content_type or "").strip().lower()
    if ct in MIME_MAP:
        return MIME_MAP[ct]
    name_lower = filename.lower()
    for ext in sorted(ALLOWED_EXTENSIONS, key=len, reverse=True):
        if name_lower.endswith(ext):
            return ext
    return ""


def _file_type_from_path(path: str | Path) -> str:
    name = str(path).lower()
    if name.endswith(".markdown"):
        return "markdown"
    return Path(name).suffix.lower().lstrip(".")


def _is_reference_path(relative_path: str) -> bool:
    return relative_path == REFERENCES_DIRNAME or relative_path.startswith(f"{REFERENCES_DIRNAME}/")


def _entry_role(relative_path: str, is_dir: bool) -> str:
    if relative_path == MEMORY_ENTRYPOINT_RELATIVE:
        return "memoryIndex"
    if _is_memory_workspace_path(relative_path):
        return "memory"
    if _is_reference_path(relative_path):
        return "reference"
    if is_dir:
        return "folder"
    return "document"


def _safe_unique_path(directory: Path, filename: str) -> Path:
    safe_name = Path(filename).name.strip() or "untitled"
    stem = Path(safe_name).stem or "untitled"
    suffix = Path(safe_name).suffix
    candidate = directory / safe_name
    index = 2
    while candidate.exists():
        candidate = directory / f"{stem} {index}{suffix}"
        index += 1
    return candidate


def _truncate_text(raw: str, *, max_lines: int, max_bytes: int) -> dict[str, Any]:
    text = raw.strip()
    lines = text.split("\n") if text else []
    line_count = len(lines)
    byte_count = len(text.encode("utf-8"))
    was_line_truncated = line_count > max_lines
    was_byte_truncated = byte_count > max_bytes
    if was_line_truncated:
        text = "\n".join(lines[:max_lines])
    encoded = text.encode("utf-8")
    if len(encoded) > max_bytes:
        text = encoded[:max_bytes].decode("utf-8", errors="ignore")
        cut_at = text.rfind("\n")
        if cut_at > 0:
            text = text[:cut_at]
    return {
        "content": text,
        "lineCount": line_count,
        "byteCount": byte_count,
        "wasLineTruncated": was_line_truncated,
        "wasByteTruncated": was_byte_truncated,
        "truncated": was_line_truncated or was_byte_truncated,
    }


def _parse_memory_frontmatter(content: str) -> dict[str, str]:
    lines = content.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    if not lines or lines[0].strip() != "---":
        return {}
    data: dict[str, str] = {}
    for line in lines[1:30]:
        if line.strip() == "---":
            break
        if ":" not in line:
            continue
        key, value = line.split(":", 1)
        key = key.strip()
        if key in {"name", "description", "type"}:
            data[key] = value.strip().strip("\"'")
    if data.get("type") and data["type"] not in MEMORY_TYPES:
        data.pop("type", None)
    return data


def _memory_body_excerpt(content: str) -> str:
    text = re.sub(r"^---\s*\n.*?\n---\s*\n?", "", content, flags=re.DOTALL).strip()
    first_line = next((line.strip() for line in text.splitlines() if line.strip()), "")
    return first_line[:140]


def _memory_index_line(memory_relative_path: str, content: str, *, title: str | None = None, hook: str | None = None) -> str:
    frontmatter = _parse_memory_frontmatter(content)
    line_title = (title or frontmatter.get("name") or Path(memory_relative_path).stem).strip()
    line_hook = (hook or frontmatter.get("description") or _memory_body_excerpt(content)).strip()
    if line_hook:
        return f"- [{line_title}]({memory_relative_path}) - {line_hook}"
    return f"- [{line_title}]({memory_relative_path})"


def _remove_memory_index_entry(entrypoint: Path, memory_relative_path: str) -> None:
    if not entrypoint.exists():
        return
    escaped = re.escape(memory_relative_path)
    pattern = re.compile(rf"^\s*-\s+\[[^\]]+\]\((?:\./)?{escaped}\)(?:\s+-.*)?\s*$")
    lines = entrypoint.read_text(encoding="utf-8", errors="replace").splitlines()
    next_lines = [line for line in lines if not pattern.match(line)]
    if next_lines != lines:
        entrypoint.write_text("\n".join(next_lines).rstrip() + ("\n" if next_lines else ""), encoding="utf-8")


def _upsert_memory_index_entry(
    workspace_id: str,
    memory_relative_path: str,
    content: str,
    *,
    title: str | None = None,
    hook: str | None = None,
) -> None:
    if memory_relative_path == MEMORY_ENTRYPOINT_FILENAME:
        return
    entrypoint = _workspace_memory_entrypoint(workspace_id)
    entrypoint.parent.mkdir(parents=True, exist_ok=True)
    if not entrypoint.exists():
        entrypoint.write_text("", encoding="utf-8")
    _remove_memory_index_entry(entrypoint, memory_relative_path)
    current = entrypoint.read_text(encoding="utf-8", errors="replace").rstrip()
    line = _memory_index_line(memory_relative_path, content, title=title, hook=hook)
    entrypoint.write_text((current + "\n" if current else "") + line + "\n", encoding="utf-8")


def _document_timestamp(path: Path) -> str:
    return datetime.fromtimestamp(path.stat().st_mtime, tz=timezone.utc).strftime("%Y-%m-%dT%H:%M:%S")


def _hash_bytes(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


def _hash_file(path: Path) -> str:
    try:
        return _hash_bytes(path.read_bytes())
    except OSError:
        return ""


def _relative_path_for_file(workspace_id: str, path: Path) -> str:
    root = _workspace_files_root(workspace_id).resolve()
    return path.resolve().relative_to(root).as_posix()


def _iter_workspace_files(workspace_id: str) -> list[tuple[str, Path]]:
    root = _workspace_files_root(workspace_id)
    files: list[tuple[str, Path]] = []
    if not root.exists():
        return files
    for path in root.rglob("*"):
        if not path.is_file():
            continue
        relative = _relative_path_for_file(workspace_id, path)
        parts = relative.split("/")
        if INTERNAL_DIRNAME in parts:
            continue
        if path.suffix.lower() not in ALLOWED_EXTENSIONS and not relative.endswith(".markdown"):
            continue
        files.append((relative, path))
    files.sort(key=lambda item: item[0].lower())
    return files


def _iter_memory_files(workspace_id: str, *, include_entrypoint: bool = True) -> list[tuple[str, str, Path]]:
    root = _workspace_memory_root(workspace_id)
    files: list[tuple[str, str, Path]] = []
    if not root.exists():
        return files
    for path in root.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in {".md", ".markdown"}:
            continue
        try:
            memory_relative = path.resolve().relative_to(root.resolve()).as_posix()
        except ValueError:
            continue
        if not include_entrypoint and memory_relative == MEMORY_ENTRYPOINT_FILENAME:
            continue
        files.append((memory_relative, _memory_workspace_path(memory_relative), path))
    files.sort(key=lambda item: item[0].lower())
    return files


def _build_tree_node(workspace_id: str, path: Path) -> dict[str, Any] | None:
    relative = _relative_path_for_file(workspace_id, path)
    if relative.split("/")[0] == INTERNAL_DIRNAME:
        return None
    is_dir = path.is_dir()
    stat = path.stat()
    ext = path.suffix.lower()
    file_type = _file_type_from_path(path)
    node: dict[str, Any] = {
        "name": path.name,
        "path": relative,
        "kind": "directory" if is_dir else "file",
        "role": _entry_role(relative, is_dir),
        "size": 0 if is_dir else stat.st_size,
        "updatedAt": _document_timestamp(path),
        "editable": (not is_dir) and ext in EDITABLE_EXTENSIONS,
        "readOnly": (not is_dir) and ext in READONLY_EXTENSIONS,
        "type": file_type,
        "extension": ext.lstrip("."),
        "isReference": _is_reference_path(relative),
    }
    if is_dir:
        children = []
        for child in sorted(path.iterdir(), key=lambda item: (not item.is_dir(), item.name.lower())):
            child_node = _build_tree_node(workspace_id, child)
            if child_node is not None:
                children.append(child_node)
        node["children"] = children
    else:
        node["contentHash"] = _hash_file(path)
    return node


def _build_memory_tree_node(workspace_id: str, path: Path) -> dict[str, Any] | None:
    memory_root = _workspace_memory_root(workspace_id).resolve()
    resolved = path.resolve()
    try:
        memory_relative = resolved.relative_to(memory_root).as_posix()
    except ValueError:
        return None
    if memory_relative == ".":
        memory_relative = ""
    is_dir = path.is_dir()
    if not is_dir and path.suffix.lower() not in {".md", ".markdown"}:
        return None
    workspace_relative = _memory_workspace_path(memory_relative)
    stat = path.stat()
    role = "memoryFolder" if is_dir else ("memoryIndex" if memory_relative == MEMORY_ENTRYPOINT_FILENAME else "memory")
    node: dict[str, Any] = {
        "name": path.name,
        "path": workspace_relative,
        "kind": "directory" if is_dir else "file",
        "role": role,
        "size": 0 if is_dir else stat.st_size,
        "updatedAt": _document_timestamp(path),
        "editable": not is_dir,
        "readOnly": False,
        "type": "" if is_dir else ("markdown" if path.suffix.lower() == ".markdown" else "md"),
        "extension": "" if is_dir else path.suffix.lower().lstrip("."),
        "isReference": False,
        "isMemory": True,
    }
    if is_dir:
        children: list[dict[str, Any]] = []
        for child in sorted(path.iterdir(), key=lambda item: (not item.is_dir(), item.name.lower())):
            child_node = _build_memory_tree_node(workspace_id, child)
            if child_node is not None:
                children.append(child_node)
        node["children"] = children
    else:
        node["contentHash"] = _hash_file(path)
    return node


def _build_openwps_tree_node(workspace_id: str) -> dict[str, Any] | None:
    memory_root = _workspace_memory_root(workspace_id)
    memory_node = _build_memory_tree_node(workspace_id, memory_root)
    if memory_node is None:
        return None
    return {
        "name": INTERNAL_DIRNAME,
        "path": INTERNAL_DIRNAME,
        "kind": "directory",
        "role": "openwps",
        "size": 0,
        "updatedAt": _document_timestamp(_workspace_internal_root(workspace_id)),
        "editable": False,
        "readOnly": True,
        "type": "",
        "extension": "",
        "isReference": False,
        "isMemory": True,
        "children": [memory_node],
    }


def _root_tree_sort_key(node: dict[str, Any]) -> tuple[int, int, str]:
    name = str(node.get("name") or "")
    path = str(node.get("path") or "")
    priority = SYSTEM_ROOT_DIR_ORDER.get(path, 100)
    is_dir = node.get("kind") == "directory"
    return (priority, 0 if is_dir else 1, name.lower())


def get_workspace_tree(workspace_id: str | None = None) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    root = _workspace_files_root(target)
    children = []
    for child in sorted(root.iterdir(), key=lambda item: (not item.is_dir(), item.name.lower())):
        node = _build_tree_node(target, child)
        if node is not None:
            children.append(node)
    openwps_node = _build_openwps_tree_node(target)
    if openwps_node is not None:
        children.append(openwps_node)
    children.sort(key=_root_tree_sort_key)
    return {
        "workspaceId": target,
        "root": {
            "name": target,
            "path": "",
            "kind": "directory",
            "role": "workspace",
            "children": children,
        },
    }


def _memory_header(memory_relative: str, workspace_relative: str, path: Path) -> dict[str, Any]:
    content = path.read_text(encoding="utf-8", errors="replace")
    frontmatter = _parse_memory_frontmatter(content)
    stat = path.stat()
    return {
        "path": workspace_relative,
        "memoryPath": memory_relative,
        "name": frontmatter.get("name") or Path(memory_relative).stem,
        "description": frontmatter.get("description") or "",
        "type": frontmatter.get("type") or "",
        "size": stat.st_size,
        "updatedAt": _document_timestamp(path),
        "contentHash": _hash_file(path),
    }


def _extract_entrypoint_links(content: str) -> dict[str, str]:
    links: dict[str, str] = {}
    for line in content.splitlines():
        for match in re.finditer(r"\[[^\]]+\]\(([^)]+)\)", line):
            target = match.group(1).strip()
            if not target or "://" in target or target.startswith("#"):
                continue
            try:
                memory_relative = _sanitize_memory_relative_path(target)
            except HTTPException:
                continue
            links[memory_relative] = line.strip()
    return links


def _query_terms(query: str | None) -> list[str]:
    text = str(query or "").lower()
    return [term for term in re.split(r"[\s,，。；;：:、/\\|()\[\]{}<>\"']+", text) if len(term) >= 2]


def _select_memory_files(
    query: str | None,
    entrypoint_content: str,
    manifest: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    terms = _query_terms(query)
    if not terms:
        return []
    entrypoint_links = _extract_entrypoint_links(entrypoint_content)
    scored: list[tuple[int, str, dict[str, Any]]] = []
    for item in manifest:
        memory_path = str(item.get("memoryPath") or "")
        haystack = " ".join([
            memory_path,
            str(item.get("name") or ""),
            str(item.get("description") or ""),
            str(item.get("type") or ""),
            entrypoint_links.get(memory_path, ""),
        ]).lower()
        score = sum(1 for term in terms if term in haystack)
        if score > 0:
            scored.append((score, str(item.get("updatedAt") or ""), item))
    scored.sort(key=lambda row: (row[0], row[1]), reverse=True)
    return [item for _score, _updated, item in scored[:MEMORY_SELECTED_LIMIT]]


def _read_selected_memory_files(workspace_id: str, selected: list[dict[str, Any]]) -> list[dict[str, Any]]:
    result: list[dict[str, Any]] = []
    for item in selected:
        memory_path = str(item.get("memoryPath") or "")
        if not memory_path:
            continue
        try:
            _memory_relative, workspace_relative, path = _resolve_memory_path(workspace_id, memory_path)
            raw = path.read_text(encoding="utf-8", errors="replace")
            limited = _truncate_text(raw, max_lines=MEMORY_FILE_MAX_LINES, max_bytes=MEMORY_FILE_MAX_BYTES)
            result.append({
                **item,
                "path": workspace_relative,
                "content": limited["content"],
                "truncated": limited["truncated"],
                "lineCount": limited["lineCount"],
                "byteCount": limited["byteCount"],
            })
        except Exception:
            continue
    return result


def get_workspace_memory(workspace_id: str | None = None, *, query: str | None = None) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    entrypoint = _workspace_memory_entrypoint(target)
    entrypoint.parent.mkdir(parents=True, exist_ok=True)
    if not entrypoint.exists():
        entrypoint.write_text("", encoding="utf-8")
    entrypoint_text = entrypoint.read_text(encoding="utf-8", errors="replace")
    entrypoint_limited = _truncate_text(
        entrypoint_text,
        max_lines=MEMORY_ENTRYPOINT_MAX_LINES,
        max_bytes=MEMORY_ENTRYPOINT_MAX_BYTES,
    )
    manifest = [
        _memory_header(memory_relative, workspace_relative, path)
        for memory_relative, workspace_relative, path in _iter_memory_files(target, include_entrypoint=False)
    ]
    manifest.sort(key=lambda item: str(item.get("updatedAt") or ""), reverse=True)
    manifest = manifest[:MEMORY_SCAN_LIMIT]
    selected_headers = _select_memory_files(query, entrypoint_text, manifest)
    return {
        "workspaceId": target,
        "entrypoint": {
            "path": MEMORY_ENTRYPOINT_RELATIVE,
            "memoryPath": MEMORY_ENTRYPOINT_FILENAME,
            "content": entrypoint_limited["content"],
            "updatedAt": _document_timestamp(entrypoint),
            "truncated": entrypoint_limited["truncated"],
            "lineCount": entrypoint_limited["lineCount"],
            "byteCount": entrypoint_limited["byteCount"],
        },
        "manifest": manifest,
        "selected": _read_selected_memory_files(target, selected_headers),
        "truncated": len(manifest) >= MEMORY_SCAN_LIMIT,
    }


def get_workspace_manifest(workspace_id: str | None = None, *, limit: int = 200, query: str | None = None) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    files = []
    references = []
    for relative, path in _iter_workspace_files(target):
        stat = path.stat()
        item = {
            "path": relative,
            "name": path.name,
            "type": _file_type_from_path(path),
            "size": stat.st_size,
            "updatedAt": _document_timestamp(path),
            "editable": path.suffix.lower() in EDITABLE_EXTENSIONS,
            "role": _entry_role(relative, False),
            "textLength": len(extract_text(path, path.suffix.lower())),
        }
        if _is_reference_path(relative):
            references.append(item)
        else:
            files.append(item)
    return {
        "workspaceId": target,
        "files": files[:limit],
        "references": references[:limit],
        "truncated": len(files) + len(references) > limit * 2,
        "memory": get_workspace_memory(target, query=query),
    }


def create_folder(workspace_id: str, path: str) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    relative, resolved = _resolve_workspace_path(target, path)
    resolved.mkdir(parents=True, exist_ok=True)
    return {"workspaceId": target, "path": relative, "kind": "directory"}


def save_file(workspace_id: str, path: str, content: bytes, *, content_type: str = "") -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    relative, resolved = _resolve_workspace_path(target, path)
    extension = _detect_extension(relative, content_type) or Path(relative).suffix.lower()
    if extension not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"不支持的文件类型。支持：{', '.join(sorted(ALLOWED_EXTENSIONS))}")
    resolved.parent.mkdir(parents=True, exist_ok=True)
    if not content and extension == ".docx":
        content = _blank_docx_bytes()
    if not content and extension in {".md", ".markdown"}:
        content = b"# Untitled\n"
    if not content and extension == ".txt":
        content = b""
    _snapshot_file(target, relative, resolved)
    resolved.write_bytes(content)
    _write_index(target, relative, resolved)
    stat = resolved.stat()
    return {
        "workspaceId": target,
        "path": relative,
        "name": resolved.name,
        "type": _file_type_from_path(resolved),
        "size": stat.st_size,
        "updatedAt": _document_timestamp(resolved),
        "contentHash": _hash_file(resolved),
        "editable": extension in EDITABLE_EXTENSIONS,
        "isReference": _is_reference_path(relative),
    }


def upload_workspace_file(workspace_id: str, directory: str | None, filename: str, content_type: str, content: bytes) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    directory_relative, directory_path = _resolve_workspace_path(target, directory or "", allow_empty=True)
    if directory_path.exists() and not directory_path.is_dir():
        raise HTTPException(status_code=400, detail="上传目标必须是目录")
    directory_path.mkdir(parents=True, exist_ok=True)
    extension = _detect_extension(filename, content_type)
    if not extension:
        raise HTTPException(status_code=400, detail=f"不支持的文件类型。支持：{', '.join(sorted(ALLOWED_EXTENSIONS))}")
    target_path = _safe_unique_path(directory_path, filename)
    relative = f"{directory_relative}/{target_path.name}" if directory_relative else target_path.name
    return save_file(target, relative, content, content_type=content_type)


def read_memory_file(workspace_id: str, path: str) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    memory_relative, workspace_relative, resolved = _resolve_memory_path(target, path)
    if not resolved.exists() or not resolved.is_file():
        raise HTTPException(status_code=404, detail="记忆文件不存在")
    content = resolved.read_text(encoding="utf-8", errors="replace")
    return {
        "workspaceId": target,
        "path": workspace_relative,
        "memoryPath": memory_relative,
        "content": content,
        "updatedAt": _document_timestamp(resolved),
        "contentHash": _hash_file(resolved),
    }


def save_memory_file(
    workspace_id: str,
    path: str,
    content: str | bytes,
    *,
    name: str | None = None,
    description: str | None = None,
    memory_type: str | None = None,
    index_title: str | None = None,
    index_hook: str | None = None,
) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    memory_relative, workspace_relative, resolved = _resolve_memory_path(target, path)
    text = content.decode("utf-8", errors="replace") if isinstance(content, bytes) else str(content)
    memory_type_normalized = memory_type if memory_type in MEMORY_TYPES else None
    if memory_relative != MEMORY_ENTRYPOINT_FILENAME and not _parse_memory_frontmatter(text) and (name or description or memory_type_normalized):
        frontmatter = [
            "---",
            f"name: {name or Path(memory_relative).stem}",
            f"description: {description or _memory_body_excerpt(text)}",
            f"type: {memory_type_normalized or 'project'}",
            "---",
            "",
        ]
        text = "\n".join(frontmatter) + text.lstrip()
    resolved.parent.mkdir(parents=True, exist_ok=True)
    _snapshot_file(target, workspace_relative, resolved)
    resolved.write_text(text, encoding="utf-8")
    _write_index(target, workspace_relative, resolved)
    if memory_relative != MEMORY_ENTRYPOINT_FILENAME:
        _upsert_memory_index_entry(
            target,
            memory_relative,
            text,
            title=index_title or name,
            hook=index_hook or description,
        )
    stat = resolved.stat()
    return {
        "workspaceId": target,
        "path": workspace_relative,
        "memoryPath": memory_relative,
        "name": resolved.name,
        "type": "markdown" if resolved.suffix.lower() == ".markdown" else "md",
        "size": stat.st_size,
        "updatedAt": _document_timestamp(resolved),
        "contentHash": _hash_file(resolved),
        "editable": True,
        "role": "memoryIndex" if memory_relative == MEMORY_ENTRYPOINT_FILENAME else "memory",
    }


def delete_memory_file(workspace_id: str, path: str) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    memory_relative, workspace_relative, resolved = _resolve_memory_path(target, path)
    if memory_relative == MEMORY_ENTRYPOINT_FILENAME:
        raise HTTPException(status_code=400, detail="MEMORY.md 是记忆入口索引，不能删除")
    if not resolved.exists() or not resolved.is_file():
        raise HTTPException(status_code=404, detail="记忆文件不存在")
    _snapshot_file(target, workspace_relative, resolved)
    resolved.unlink()
    _remove_memory_index_entry(_workspace_memory_entrypoint(target), memory_relative)
    return {"success": True, "workspaceId": target, "path": workspace_relative, "memoryPath": memory_relative}


def move_memory_file(workspace_id: str, path: str, to_path: str) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    memory_relative, workspace_relative, resolved = _resolve_memory_path(target, path)
    next_memory_relative, next_workspace_relative, next_resolved = _resolve_memory_path(target, to_path)
    if memory_relative == MEMORY_ENTRYPOINT_FILENAME:
        raise HTTPException(status_code=400, detail="MEMORY.md 是固定入口索引，不能移动或重命名")
    if next_memory_relative == MEMORY_ENTRYPOINT_FILENAME:
        raise HTTPException(status_code=400, detail="不能覆盖 MEMORY.md")
    if not resolved.exists() or not resolved.is_file():
        raise HTTPException(status_code=404, detail="记忆文件不存在")
    if next_resolved.exists():
        raise HTTPException(status_code=409, detail="目标记忆文件已存在")
    next_resolved.parent.mkdir(parents=True, exist_ok=True)
    _snapshot_file(target, workspace_relative, resolved)
    resolved.rename(next_resolved)
    _remove_memory_index_entry(_workspace_memory_entrypoint(target), memory_relative)
    content = next_resolved.read_text(encoding="utf-8", errors="replace")
    _upsert_memory_index_entry(target, next_memory_relative, content)
    _write_index(target, next_workspace_relative, next_resolved)
    return {"success": True, "workspaceId": target, "path": next_workspace_relative, "memoryPath": next_memory_relative}


def read_file_path(workspace_id: str, path: str) -> Path:
    target = _require_workspace(workspace_id)
    _, resolved = _resolve_workspace_path(target, path)
    if not resolved.exists() or not resolved.is_file():
        raise HTTPException(status_code=404, detail="文件不存在")
    return resolved


def delete_file(workspace_id: str, path: str) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    relative, resolved = _resolve_workspace_path(target, path)
    if not resolved.exists():
        raise HTTPException(status_code=404, detail="文件不存在")
    if resolved.is_dir():
        shutil.rmtree(resolved)
    else:
        _snapshot_file(target, relative, resolved)
        resolved.unlink()
    return {"success": True, "workspaceId": target, "path": relative}


def move_file(workspace_id: str, path: str, to_path: str) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    relative, resolved = _resolve_workspace_path(target, path)
    next_relative, next_resolved = _resolve_workspace_path(target, to_path)
    if not resolved.exists():
        raise HTTPException(status_code=404, detail="文件不存在")
    if next_resolved.exists():
        raise HTTPException(status_code=409, detail="目标路径已存在")
    next_resolved.parent.mkdir(parents=True, exist_ok=True)
    _snapshot_file(target, relative, resolved)
    resolved.rename(next_resolved)
    if next_resolved.is_file():
        _write_index(target, next_relative, next_resolved)
    return {"success": True, "workspaceId": target, "path": next_relative}


def _index_path(workspace_id: str, relative_path: str) -> Path:
    digest = hashlib.sha256(relative_path.encode("utf-8")).hexdigest()
    return _workspace_internal_root(workspace_id) / "index" / f"{digest}.txt"


def _write_index(workspace_id: str, relative_path: str, file_path: Path) -> None:
    try:
        text = extract_text(file_path, file_path.suffix.lower())
        _index_path(workspace_id, relative_path).write_text(text, encoding="utf-8")
    except Exception as exc:
        logger.warning("Failed to index workspace file %s: %s", relative_path, exc)


def _read_index_or_extract(workspace_id: str, relative_path: str, file_path: Path) -> str:
    index_path = _index_path(workspace_id, relative_path)
    if index_path.exists() and index_path.stat().st_mtime >= file_path.stat().st_mtime:
        return index_path.read_text(encoding="utf-8", errors="replace")
    text = extract_text(file_path, file_path.suffix.lower())
    index_path.parent.mkdir(parents=True, exist_ok=True)
    index_path.write_text(text, encoding="utf-8")
    return text


def _snapshot_file(workspace_id: str, relative_path: str, file_path: Path) -> None:
    if not file_path.exists() or not file_path.is_file():
        return
    version_dir = _workspace_internal_root(workspace_id) / "versions"
    version_dir.mkdir(parents=True, exist_ok=True)
    safe_name = re.sub(r"[^A-Za-z0-9_.-]+", "_", relative_path.replace("/", "__"))
    snapshot = version_dir / f"{datetime.now(timezone.utc).strftime('%Y%m%dT%H%M%SZ')}_{safe_name}"
    try:
        shutil.copy2(file_path, snapshot)
    except Exception as exc:
        logger.warning("Failed to snapshot workspace file %s: %s", relative_path, exc)


def open_file_as_document(workspace_id: str, path: str) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    relative, resolved = _resolve_workspace_path(target, path, allow_internal=_is_memory_workspace_path(path))
    if not resolved.exists() or not resolved.is_file():
        raise HTTPException(status_code=404, detail="文件不存在")
    extension = resolved.suffix.lower()
    if extension not in EDITABLE_EXTENSIONS:
        raise HTTPException(status_code=400, detail="该文件类型暂不支持直接编辑")
    doc_json, page_config, extra = _file_to_doc_payload(resolved, extension)
    stat = resolved.stat()
    payload = {
        "workspaceId": target,
        "filePath": relative,
        "fileType": _file_type_from_path(resolved),
        "fileMtime": stat.st_mtime,
        "contentHash": _hash_file(resolved),
        "currentDocumentName": resolved.name,
        "docJson": doc_json,
        "pageConfig": page_config,
        "dirty": False,
    }
    payload.update(extra)
    return payload


def save_document_session_file(
    workspace_id: str,
    path: str,
    doc_json: dict[str, Any],
    page_config: dict[str, Any] | None = None,
) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    if _is_memory_workspace_path(path):
        memory_relative, _workspace_relative, resolved = _resolve_memory_path(target, path)
        extension = resolved.suffix.lower()
        if extension not in {".md", ".markdown"}:
            raise HTTPException(status_code=400, detail="记忆文件只支持 Markdown")
        content = _doc_json_to_file_bytes(doc_json, extension, page_config or DEFAULT_PAGE_CONFIG)
        return save_memory_file(target, memory_relative, content)
    relative, resolved = _resolve_workspace_path(target, path)
    extension = resolved.suffix.lower()
    if extension not in EDITABLE_EXTENSIONS:
        raise HTTPException(status_code=400, detail="该文件类型暂不支持直接编辑")
    content = _doc_json_to_file_bytes(doc_json, extension, page_config or DEFAULT_PAGE_CONFIG)
    return save_file(target, relative, content)


def _blank_docx_bytes() -> bytes:
    try:
        from docx import Document
        buf = io.BytesIO()
        Document().save(buf)
        return buf.getvalue()
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"创建空 DOCX 失败：{exc}") from exc


def _node_io_worker_path() -> Path:
    return BASE_DIR / "node" / ".generated" / "server" / "node" / "document_io_worker.js"


def _node_io_file_type(extension: str) -> str:
    return extension.lower().lstrip(".")


def _run_node_io_worker(payload: dict[str, Any]) -> dict[str, Any] | None:
    worker = _node_io_worker_path()
    if not worker.exists():
        return None
    try:
        proc = subprocess.run(
            ["node", str(worker)],
            input=json.dumps(payload, ensure_ascii=False),
            capture_output=True,
            text=True,
            cwd=str(BASE_DIR.parent),
            timeout=90,
            check=False,
        )
    except Exception as exc:
        logger.warning("Failed to run document IO worker: %s", exc)
        return None
    if proc.returncode != 0:
        logger.warning("Document IO worker exited with %s: %s", proc.returncode, proc.stderr.strip())
        return None
    try:
        result = json.loads(proc.stdout or "{}")
    except json.JSONDecodeError as exc:
        logger.warning("Document IO worker returned invalid JSON: %s", exc)
        return None
    if not isinstance(result, dict) or result.get("success") is not True:
        message = result.get("message") if isinstance(result, dict) else result
        logger.warning("Document IO worker failed: %s", message)
        return None
    return result


def _open_via_node_io(path: Path, extension: str) -> dict[str, Any] | None:
    if extension not in EDITABLE_EXTENSIONS:
        return None
    return _run_node_io_worker({
        "operation": "open",
        "name": path.name,
        "fileType": _node_io_file_type(extension),
        "contentBase64": base64.b64encode(path.read_bytes()).decode("ascii"),
    })


def _save_via_node_io(
    doc_json: dict[str, Any],
    extension: str,
    page_config: dict[str, Any],
    export_options: dict[str, Any] | None = None,
) -> bytes | None:
    if extension not in EDITABLE_EXTENSIONS:
        return None
    result = _run_node_io_worker({
        "operation": "save",
        "fileType": _node_io_file_type(extension),
        "docJson": doc_json,
        "pageConfig": page_config,
        "exportOptions": export_options or {},
    })
    content_base64 = result.get("contentBase64") if isinstance(result, dict) else None
    if not isinstance(content_base64, str):
        return None
    try:
        return base64.b64decode(content_base64)
    except Exception as exc:
        logger.warning("Document IO worker returned invalid base64: %s", exc)
        return None


def _paragraph_node(text: str, attrs: dict[str, Any] | None = None) -> dict[str, Any]:
    node = {"type": "paragraph", "attrs": {**DEFAULT_PARAGRAPH_ATTRS, **(attrs or {})}}
    if text:
        node["content"] = [{"type": "text", "text": text}]
    return node


def _text_to_doc_json(text: str, *, markdown: bool = False) -> dict[str, Any]:
    lines = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    content: list[dict[str, Any]] = []
    for line in lines:
        attrs: dict[str, Any] = {}
        value = line
        if markdown:
            heading = re.match(r"^(#{1,6})\s+(.*)$", line)
            bullet = re.match(r"^\s*[-*+]\s+(.*)$", line)
            ordered = re.match(r"^\s*\d+[.)]\s+(.*)$", line)
            if heading:
                attrs["headingLevel"] = len(heading.group(1))
                value = heading.group(2)
            elif bullet:
                attrs["listType"] = "bullet"
                value = bullet.group(1)
            elif ordered:
                attrs["listType"] = "ordered"
                value = ordered.group(1)
        content.append(_paragraph_node(value, attrs))
    return {"type": "doc", "content": content or [_paragraph_node("")]}


def _docx_to_doc_json(path: Path) -> tuple[dict[str, Any], dict[str, Any]]:
    try:
        from docx import Document
        document = Document(str(path))
        content: list[dict[str, Any]] = []
        for paragraph in document.paragraphs:
            text = paragraph.text or ""
            attrs: dict[str, Any] = {}
            style_name = (paragraph.style.name if paragraph.style else "") or ""
            match = re.search(r"heading\s+([1-6])", style_name, re.IGNORECASE)
            if match:
                attrs["headingLevel"] = int(match.group(1))
            content.append(_paragraph_node(text, attrs))
        for table in document.tables:
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cell_text = "\n".join(p.text for p in cell.paragraphs).strip()
                    cells.append({
                        "type": "table_cell",
                        "attrs": {
                            "header": False,
                            "colspan": 1,
                            "rowspan": 1,
                            "width": None,
                            "backgroundColor": "",
                            "borderColor": "#cccccc",
                            "borderWidth": 1,
                        },
                        "content": [_paragraph_node(cell_text)],
                    })
                rows.append({"type": "table_row", "content": cells})
            if rows:
                content.append({"type": "table", "content": rows})
        return {"type": "doc", "content": content or [_paragraph_node("")]}, dict(DEFAULT_PAGE_CONFIG)
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"DOCX 打开失败：{exc}") from exc


def _file_to_doc_json_fallback(path: Path, extension: str) -> tuple[dict[str, Any], dict[str, Any]]:
    if extension == ".docx":
        return _docx_to_doc_json(path)
    text = path.read_text(encoding="utf-8", errors="replace")
    return _text_to_doc_json(text, markdown=extension in {".md", ".markdown"}), dict(DEFAULT_PAGE_CONFIG)


def _file_to_doc_payload(path: Path, extension: str) -> tuple[dict[str, Any], dict[str, Any], dict[str, Any]]:
    node_result = _open_via_node_io(path, extension)
    if node_result:
        doc_json = node_result.get("docJson")
        page_config = node_result.get("pageConfig")
        if isinstance(doc_json, dict) and isinstance(page_config, dict):
            extra: dict[str, Any] = {}
            export_options = node_result.get("exportOptions")
            if isinstance(export_options, dict):
                extra["docxExportOptions"] = export_options
                typography = export_options.get("typography")
                if isinstance(typography, dict) and typography.get("punctuationCompression") is True:
                    extra["docxLetterSpacingPx"] = -0.34
            return doc_json, page_config, extra
    doc_json, page_config = _file_to_doc_json_fallback(path, extension)
    return doc_json, page_config, {}


def _file_to_doc_json(path: Path, extension: str) -> tuple[dict[str, Any], dict[str, Any]]:
    doc_json, page_config, _extra = _file_to_doc_payload(path, extension)
    return doc_json, page_config


def _node_text(node: Any) -> str:
    if not isinstance(node, dict):
        return ""
    text = node.get("text")
    if isinstance(text, str):
        return text
    content = node.get("content")
    if not isinstance(content, list):
        return ""
    return "".join(_node_text(item) for item in content)


def _doc_json_to_plain_lines(doc_json: dict[str, Any]) -> list[str]:
    lines: list[str] = []
    for node in doc_json.get("content") or []:
        if not isinstance(node, dict):
            continue
        node_type = node.get("type")
        if node_type == "paragraph":
            lines.append(_node_text(node))
        elif node_type == "table":
            for row in node.get("content") or []:
                cells = [_node_text(cell) for cell in row.get("content") or [] if isinstance(cell, dict)]
                lines.append(" | ".join(cells))
        elif node_type == "horizontal_rule":
            lines.append("---")
    return lines


def _doc_json_to_markdown(doc_json: dict[str, Any]) -> str:
    lines: list[str] = []
    for node in doc_json.get("content") or []:
        if not isinstance(node, dict):
            continue
        if node.get("type") == "paragraph":
            attrs = node.get("attrs") if isinstance(node.get("attrs"), dict) else {}
            text = _node_text(node)
            heading = attrs.get("headingLevel")
            list_type = attrs.get("listType")
            if isinstance(heading, int) and 1 <= heading <= 6:
                lines.append(f"{'#' * heading} {text}".rstrip())
            elif list_type == "bullet":
                lines.append(f"- {text}".rstrip())
            elif list_type == "ordered":
                lines.append(f"1. {text}".rstrip())
            else:
                lines.append(text)
        elif node.get("type") == "table":
            table_rows: list[list[str]] = []
            for row in node.get("content") or []:
                if not isinstance(row, dict):
                    continue
                table_rows.append([_node_text(cell) for cell in row.get("content") or [] if isinstance(cell, dict)])
            if table_rows:
                lines.append("| " + " | ".join(table_rows[0]) + " |")
                lines.append("| " + " | ".join("---" for _ in table_rows[0]) + " |")
                for row in table_rows[1:]:
                    lines.append("| " + " | ".join(row) + " |")
        elif node.get("type") == "horizontal_rule":
            lines.append("---")
    return "\n".join(lines).rstrip() + "\n"


def _doc_json_to_docx_bytes(doc_json: dict[str, Any]) -> bytes:
    try:
        from docx import Document
        document = Document()
        for node in doc_json.get("content") or []:
            if not isinstance(node, dict):
                continue
            if node.get("type") == "paragraph":
                attrs = node.get("attrs") if isinstance(node.get("attrs"), dict) else {}
                text = _node_text(node)
                heading = attrs.get("headingLevel")
                if isinstance(heading, int) and 1 <= heading <= 6:
                    document.add_heading(text, level=heading)
                else:
                    paragraph = document.add_paragraph()
                    list_type = attrs.get("listType")
                    if list_type == "bullet":
                        paragraph.style = "List Bullet"
                    elif list_type == "ordered":
                        paragraph.style = "List Number"
                    paragraph.add_run(text)
            elif node.get("type") == "table":
                rows = node.get("content") if isinstance(node.get("content"), list) else []
                if not rows:
                    continue
                max_cols = max((len(row.get("content") or []) for row in rows if isinstance(row, dict)), default=0)
                if max_cols <= 0:
                    continue
                table = document.add_table(rows=len(rows), cols=max_cols)
                for row_index, row in enumerate(rows):
                    if not isinstance(row, dict):
                        continue
                    for col_index, cell in enumerate(row.get("content") or []):
                        if col_index < max_cols:
                            table.cell(row_index, col_index).text = _node_text(cell)
            elif node.get("type") == "horizontal_rule":
                document.add_paragraph("---")
        buf = io.BytesIO()
        document.save(buf)
        return buf.getvalue()
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"DOCX 保存失败：{exc}") from exc


def _doc_json_to_file_bytes(doc_json: dict[str, Any], extension: str, page_config: dict[str, Any]) -> bytes:
    node_content = _save_via_node_io(doc_json, extension, page_config)
    if node_content is not None:
        return node_content
    if extension == ".docx":
        return _doc_json_to_docx_bytes(doc_json)
    if extension in {".md", ".markdown"}:
        return _doc_json_to_markdown(doc_json).encode("utf-8")
    return ("\n".join(_doc_json_to_plain_lines(doc_json))).encode("utf-8")


def _render_pdf_pages_to_images(file_path: Path, max_pages: int = 10) -> list[tuple[int, str]]:
    try:
        import pypdfium2
    except ImportError:
        return []

    pages_data: list[tuple[int, str]] = []
    try:
        pdf = pypdfium2.PdfDocument(str(file_path))
        page_count = min(len(pdf), max_pages)
        for i in range(page_count):
            page = pdf[i]
            bitmap = page.render(scale=2)
            img = bitmap.to_pil()
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            b64 = base64.b64encode(buf.getvalue()).decode("ascii")
            pages_data.append((i + 1, f"data:image/png;base64,{b64}"))
        pdf.close()
    except Exception as exc:
        logger.warning("Failed to render PDF pages: %s", exc)
    return pages_data


PADDLEOCR_API_URL = "https://60acw0te7fv2z9yf.aistudio-app.com/layout-parsing"
PADDLEOCR_TOKEN = "baf2212551eec93343531e64ce30a103e8fcf934"


def _ocr_pdf_via_paddleocr(file_path: Path) -> str:
    try:
        file_bytes = file_path.read_bytes()
        file_data = base64.b64encode(file_bytes).decode("ascii")
        payload = {
            "file": file_data,
            "fileType": 0,
            "useDocOrientationClassify": False,
            "useDocUnwarping": False,
            "useChartRecognition": False,
        }
        resp = httpx.post(
            PADDLEOCR_API_URL,
            json=payload,
            headers={
                "Authorization": f"token {PADDLEOCR_TOKEN}",
                "Content-Type": "application/json",
            },
            timeout=120,
        )
        if resp.status_code != 200:
            logger.warning("PaddleOCR API returned status %d: %s", resp.status_code, resp.text[:500])
            return ""
        data = resp.json()
        result = data.get("result", {})
        layout_results = result.get("layoutParsingResults", [])
        pages_text: list[str] = []
        for i, res in enumerate(layout_results):
            md = res.get("markdown", {})
            text = (md.get("text") or "").strip()
            if text:
                pages_text.append(f"--- 第 {i + 1} 页 ---\n{text}")
        return "\n\n".join(pages_text)
    except Exception as exc:
        logger.warning("PaddleOCR layout parsing failed for PDF: %s", exc)
        return ""


def extract_text(file_path: Path, extension: str) -> str:
    try:
        if extension in {".txt", ".md", ".markdown"}:
            return file_path.read_text(encoding="utf-8", errors="replace")
        if extension == ".docx":
            from docx import Document
            doc = Document(str(file_path))
            parts: list[str] = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            return "\n".join(parts)
        if extension == ".pdf":
            text = _extract_pdf_text(file_path)
            if text.strip():
                return text
            ocr_text = _ocr_pdf_via_paddleocr(file_path)
            if ocr_text.strip():
                return ocr_text
            page_count = _pdf_page_count(file_path)
            return f"[扫描版PDF，共{page_count}页，文字提取为空。此文件可能是扫描件或图片PDF，OCR服务未配置或未能识别文字。]"
        if extension in {".ppt", ".pptx"}:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides: list[str] = []
            for i, slide in enumerate(prs.slides):
                texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                if texts:
                    slides.append(f"--- 幻灯片 {i + 1} ---\n" + "\n".join(texts))
            return "\n\n".join(slides)
    except Exception as exc:
        return f"[提取失败: {exc}]"
    return "[不支持的文件格式]"


def _extract_pdf_text(file_path: Path) -> str:
    try:
        import pdfplumber
        pages: list[str] = []
        with pdfplumber.open(str(file_path)) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"--- 第 {i + 1} 页 ---\n{text}")
        return "\n\n".join(pages)
    except ImportError:
        pass
    except Exception:
        pass
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(file_path))
        pages: list[str] = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"--- 第 {i + 1} 页 ---\n{text}")
        return "\n\n".join(pages)
    except Exception:
        return ""


def _pdf_page_count(file_path: Path) -> int:
    try:
        import pdfplumber
        with pdfplumber.open(str(file_path)) as pdf:
            return len(pdf.pages)
    except Exception:
        pass
    try:
        from PyPDF2 import PdfReader
        return len(PdfReader(str(file_path)).pages)
    except Exception:
        return 0


def search_text_lines(text: str, query: str, context_lines: int = 3) -> list[dict[str, Any]]:
    keywords = [kw for kw in query.lower().split() if kw]
    if not keywords:
        return []
    lines = text.split("\n")
    results: list[dict[str, Any]] = []
    for i, line in enumerate(lines):
        line_lower = line.lower()
        if all(kw in line_lower for kw in keywords):
            start = max(0, i - context_lines)
            end = min(len(lines), i + context_lines + 1)
            results.append({
                "lineNumber": i,
                "matchedLine": line,
                "context": lines[start:end],
                "contextStart": start,
            })
    return results


def _scope_accepts_path(relative: str, *, scope: str, path: str | None) -> bool:
    if path:
        normalized = _sanitize_relative_path(path, allow_empty=True, allow_internal=_is_memory_workspace_path(path))
        if normalized and not (relative == normalized or relative.startswith(f"{normalized}/")):
            return False
    if scope == "references":
        return _is_reference_path(relative)
    if scope == "workspace":
        return not _is_reference_path(relative)
    if scope == "memory":
        return _is_memory_workspace_path(relative)
    return True


def search_workspace(
    query: str,
    doc_id: str | None = None,
    context_lines: int = 3,
    *,
    workspace_id: str | None = None,
    scope: str = "all",
    path: str | None = None,
) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    scope_normalized = str(scope or "all").strip().lower()
    if scope_normalized not in {"all", "workspace", "references", "path", "memory"}:
        scope_normalized = "all"
    path_filter = path or doc_id
    all_results: list[dict[str, Any]] = []
    total = 0
    source_files: list[tuple[str, Path]]
    if scope_normalized == "memory" or (path_filter and _is_memory_workspace_path(path_filter)):
        source_files = [(workspace_relative, file_path) for _memory_relative, workspace_relative, file_path in _iter_memory_files(target)]
    else:
        source_files = _iter_workspace_files(target)
    for relative, file_path in source_files:
        if not _scope_accepts_path(relative, scope=scope_normalized, path=path_filter):
            continue
        total += 1
        text = _read_index_or_extract(target, relative, file_path)
        matches = search_text_lines(text, query, context_lines)
        if matches:
            all_results.append({
                "docId": relative,
                "docPath": relative,
                "docName": file_path.name,
                "matches": matches,
                "matchCount": len(matches),
            })
    return {
        "query": query,
        "workspaceId": target,
        "scope": scope_normalized,
        "path": path_filter,
        "totalDocs": total,
        "matchedDocs": len(all_results),
        "results": all_results,
    }


def get_document_content(
    doc_id: str,
    from_line: int | None = None,
    to_line: int | None = None,
    *,
    workspace_id: str | None = None,
) -> dict[str, Any]:
    target = _require_workspace(workspace_id)
    relative, file_path = _resolve_workspace_path(target, doc_id, allow_internal=_is_memory_workspace_path(doc_id))
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="文档不存在")
    text = _read_index_or_extract(target, relative, file_path)
    lines = text.split("\n")
    total_lines = len(lines)
    start = max(0, from_line if from_line is not None else 0)
    end = min(total_lines, to_line if to_line is not None else total_lines)
    return {
        "id": relative,
        "path": relative,
        "name": file_path.name,
        "type": _file_type_from_path(file_path),
        "workspaceId": target,
        "totalLines": total_lines,
        "fromLine": start,
        "toLine": end,
        "content": "\n".join(lines[start:end]),
    }


def list_workspace_docs() -> list[dict[str, Any]]:
    target = _require_workspace(None)
    docs: list[dict[str, Any]] = []
    for relative, path in _iter_workspace_files(target):
        if not _is_reference_path(relative):
            continue
        stat = path.stat()
        text = _read_index_or_extract(target, relative, path)
        docs.append({
            "id": relative,
            "path": relative,
            "name": path.name,
            "type": _file_type_from_path(path),
            "size": stat.st_size,
            "textLength": len(text),
            "uploadedAt": _document_timestamp(path),
            "role": "reference",
            "workspaceId": target,
        })
    return docs
